import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from mcp_message import MCPMessage

def parse_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def parse_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(path):
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    return "\n".join(slides_text)

def parse_csv(path):
    df = pd.read_csv(path)
    return df.to_string()

def parse_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def IngestionAgent_MCP(message):
    print(f"[IngestionAgent] Received: {message}")
    chunks = []
    for file in message.payload['files']:
        ext = file['name'].split('.')[-1].lower()
        try:
            if ext == 'pdf':
                text = parse_pdf(file['path'])
            elif ext == 'docx':
                text = parse_docx(file['path'])
            elif ext == 'pptx':
                text = parse_pptx(file['path'])
            elif ext == 'csv':
                text = parse_csv(file['path'])
            elif ext in ['txt', 'md']:
                text = parse_txt(file['path'])
            else:
                text = ""
            if text:
                for chunk in text.split("\n\n"):
                    if chunk.strip():
                        chunks.append({"text": chunk.strip(), "source": file['name']})
        except Exception as e:
            print(f"Error parsing {file['name']}: {e}")

    return message.__class__(
        sender="IngestionAgent",
        receiver="CoordinatorAgent",
        type_="INGESTION_RESULT",
        trace_id=message.trace_id,
        payload={"chunks": chunks}
    )